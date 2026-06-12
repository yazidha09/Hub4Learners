from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG    = RGBColor(0xF1, 0xF5, 0xF9)
DARK       = RGBColor(0x1E, 0x29, 0x3B)
SUBTLE     = RGBColor(0x47, 0x52, 0x66)
MUTED      = RGBColor(0x94, 0xA3, 0xB8)
BORDER     = RGBColor(0xE2, 0xE8, 0xF0)
ACCENT     = RGBColor(0x25, 0x63, 0xEB)
ACCENT2    = RGBColor(0xF9, 0x73, 0x16)
GREEN      = RGBColor(0x05, 0x9C, 0x69)
YELLOW     = RGBColor(0xD9, 0x77, 0x06)
PURPLE     = RGBColor(0x7C, 0x3A, 0xED)
TEAL       = RGBColor(0x0D, 0x94, 0x8A)
RED        = RGBColor(0xDC, 0x26, 0x26)

sprint_colors = [ACCENT, ACCENT2, GREEN, PURPLE, TEAL, RED]

logo_path = "frontend/src/assets/logo/icon.PNG"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(tf, text, font_size=16, color=DARK, bold=False, space_before=Pt(2), space_after=Pt(1), font_name='Calibri', alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    p.alignment = alignment
    return p


def add_page_number(slide, num, total=6):
    add_textbox(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                f"{num}/{total}", font_size=10, color=MUTED, alignment=PP_ALIGN.RIGHT)


# ─── Sprint data ─────────────────────────────────────────────────────────────
sprints = [
    {
        "id": "Sprint 1",
        "title": "Identity Foundation",
        "duration": "2 weeks",
        "subtitle": "Authentication, Roles & User Management",
        "items": [
            "User registration & login (JWT-based)",
            "Role hierarchy: Student, Professor, Admin, Super Admin",
            "Profile management (bio, avatar, settings)",
            "Base layout & responsive navigation",
            "Home page with course catalog",
            "Student & Professor dashboards",
        ],
        "icon": "🔐",
    },
    {
        "id": "Sprint 2",
        "title": "Course Architecture",
        "duration": "2 weeks",
        "subtitle": "Content Creation & Enrollment",
        "items": [
            "Course CRUD (title, description, thumbnail, category)",
            "Section / Subsection / Lesson block hierarchy",
            "Rich text editor (TipTap integration)",
            "File uploads (images, videos, PDFs)",
            "Free enrollment & Stripe checkout for paid courses",
            "Course feedback & ratings",
        ],
        "icon": "📚",
    },
    {
        "id": "Sprint 3",
        "title": "AI Learning Features",
        "duration": "2 weeks",
        "subtitle": "RAG Tutor, QCM & PDF Course Generation",
        "items": [
            "PDF-to-course generation (PyMuPDF + Gemini)",
            "RAG pipeline (Pinecone vector store + Gemini embeddings)",
            "In-course AI tutor (contextual Q&A)",
            "QCM quiz generation per subsection",
            "Discussion threads per subsection",
            "AI-powered discussion summaries",
        ],
        "icon": "🤖",
    },
    {
        "id": "Sprint 4",
        "title": "Gamification & Socials",
        "duration": "2 weeks",
        "subtitle": "XP, Levels, Streaks, Achievements & Leaderboards",
        "items": [
            "XP logging system (auditable, anti-cheat)",
            "Leveling curve (deterministic progression)",
            "Daily streaks & engagement tracking",
            "Achievement & badge system (event-driven)",
            "Leaderboards (daily, weekly, all-time)",
            "Friends list & social profiles",
        ],
        "icon": "🏆",
    },
    {
        "id": "Sprint 5",
        "title": "Communication & Community",
        "duration": "2 weeks",
        "subtitle": "Real-Time Messaging, Notifications & Discussions",
        "items": [
            "Real-time chat (WebSocket with auto-reconnect)",
            "Per-lesson discussion forums (posts, replies, votes)",
            "Content reporting system",
            "Live notifications (in-app, real-time)",
            "University announcements",
            "Online presence indicators",
        ],
        "icon": "💬",
    },
    {
        "id": "Sprint 6",
        "title": "Payment, Analytics & Admin",
        "duration": "2 weeks",
        "subtitle": "Monetization, Insights & Platform Management",
        "items": [
            "Stripe Checkout integration & webhooks",
            "Pro subscription ($9.99/month)",
            "Learner analytics (progress, time spent, scores)",
            "Professor analytics (course engagement, revenue)",
            "Multi-tenant admin dashboards",
            "Performance optimization & final polish",
        ],
        "icon": "📊",
    },
]

TOTAL = len(sprints)

for i, sprint in enumerate(sprints):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    color = sprint_colors[i]

    # Thin accent bar at top
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), color)

    # Logo
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.8), Inches(0.3), Inches(0.55), Inches(0.55))

    # Sprint badge
    badge = add_rounded_rect(slide, Inches(1.6), Inches(0.3), Inches(1.6), Inches(0.45), color)
    add_textbox(slide, Inches(1.6), Inches(0.33), Inches(1.6), Inches(0.4),
                sprint["id"], font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Title
    add_textbox(slide, Inches(3.4), Inches(0.25), Inches(8), Inches(0.5),
                sprint["title"], font_size=26, color=DARK, bold=True)

    # Subtitle + duration
    add_textbox(slide, Inches(3.4), Inches(0.75), Inches(8), Inches(0.3),
                f'{sprint["subtitle"]}  |  Duration: {sprint["duration"]}', font_size=13, color=SUBTLE)

    # Separator line
    add_shape(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.02), BORDER)

    # ─── Left panel: Sprint description card ─────────────────────────────
    card = add_rounded_rect(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.8), CARD_BG, BORDER)
    add_shape(slide, Inches(0.8), Inches(1.6), Inches(0.06), Inches(4.8), color)

    # Section header in card
    add_textbox(slide, Inches(1.2), Inches(1.8), Inches(4.8), Inches(0.35),
                "Key Deliverables", font_size=16, color=color, bold=True)

    txBox = add_textbox(slide, Inches(1.2), Inches(2.3), Inches(4.8), Inches(3.8),
                        "", font_size=12, color=DARK)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = sprint["items"][0]
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = DARK
    tf.paragraphs[0].font.name = 'Calibri'
    for item in sprint["items"][1:]:
        add_paragraph(tf, item, font_size=12, color=DARK, space_before=Pt(6))

    # ─── Right panel: Visual / Summary ────────────────────────────────────
    # Icon circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(1.8), Inches(1.2), Inches(1.2))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    add_textbox(slide, Inches(7.5), Inches(1.9), Inches(1.2), Inches(1.0),
                sprint["icon"], font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Quick summary card
    add_rounded_rect(slide, Inches(9.2), Inches(1.8), Inches(3.2), Inches(1.2), color)
    add_textbox(slide, Inches(9.4), Inches(1.9), Inches(2.8), Inches(0.3),
                "Sprint Goal", font_size=11, color=WHITE, bold=True)
    add_textbox(slide, Inches(9.4), Inches(2.2), Inches(2.8), Inches(0.6),
                f"Establish the {sprint['title'].lower()} layer of the Hub4Learners platform",
                font_size=11, color=WHITE)

    # Feature count
    add_rounded_rect(slide, Inches(7.5), Inches(3.4), Inches(2.2), Inches(1.2), CARD_BG, BORDER)
    add_textbox(slide, Inches(7.6), Inches(3.5), Inches(2.0), Inches(0.3),
                "Deliverables", font_size=11, color=MUTED, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.6), Inches(3.8), Inches(2.0), Inches(0.5),
                str(len(sprint["items"])), font_size=32, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    # Duration card
    add_rounded_rect(slide, Inches(10.0), Inches(3.4), Inches(2.4), Inches(1.2), CARD_BG, BORDER)
    add_textbox(slide, Inches(10.1), Inches(3.5), Inches(2.2), Inches(0.3),
                "Duration", font_size=11, color=MUTED, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(10.1), Inches(3.8), Inches(2.2), Inches(0.5),
                sprint["duration"], font_size=32, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    # ─── Bottom methodology bar ───────────────────────────────────────────
    add_rounded_rect(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5), CARD_BG, BORDER)
    add_shape(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.02), color)
    add_textbox(slide, Inches(1.0), Inches(6.75), Inches(11.3), Inches(0.4),
                f"Sprint {i+1} of {TOTAL}  |  Adapted Scrum — 2-week cycles, backlog refinement, demo & retrospective",
                font_size=10, color=MUTED, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, i + 1, TOTAL)

# ─── Save ──────────────────────────────────────────────────────────────────
output_path = "Sprint_Identity.pptx"
prs.save(output_path)
print(f"[OK] Sprint identity presentation saved: {output_path}")
