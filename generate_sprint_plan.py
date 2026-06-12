from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT   = RGBColor(0xF8, 0xFA, 0xFC)
CARD_BG    = RGBColor(0xF1, 0xF5, 0xF9)
DARK       = RGBColor(0x1E, 0x29, 0x3B)
SUBTLE     = RGBColor(0x47, 0x52, 0x66)
MUTED      = RGBColor(0x94, 0xA3, 0xB8)
BORDER     = RGBColor(0xE2, 0xE8, 0xF0)
ACCENT     = RGBColor(0x25, 0x63, 0xEB)
ACCENT2    = RGBColor(0xF9, 0x73, 0x16)
GREEN      = RGBColor(0x05, 0x9C, 0x69)

logo_path = "frontend/src/assets/logo/icon.PNG"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_page_number(slide, num, total=1):
    add_textbox(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                f"{num}/{total}", font_size=10, color=MUTED, alignment=PP_ALIGN.RIGHT)


def add_table(slide, rows, cols, left, top, width, height):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    return table


def set_cell(table, row, col, text, font_size=11, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, fill_color=None):
    cell = table.cell(row, col)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.TOP
    if fill_color:
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '{:02X}{:02X}{:02X}'.format(*fill_color)})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)
    # Set margins
    tcPr = cell._tc.get_or_add_tcPr()
    tcPr.set('marL', str(Emu(Inches(0.1))))
    tcPr.set('marR', str(Emu(Inches(0.1))))
    tcPr.set('marT', str(Emu(Inches(0.05))))
    tcPr.set('marB', str(Emu(Inches(0.05))))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Sprint Plan
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

# Top accent bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)

# Logo
if os.path.exists(logo_path):
    slide.shapes.add_picture(logo_path, Inches(0.8), Inches(0.3), Inches(0.6), Inches(0.6))

# Title
add_textbox(slide, Inches(1.6), Inches(0.3), Inches(10), Inches(0.5),
            "Hub4Learners — Sprint Plan", font_size=28, color=DARK, bold=True)
add_textbox(slide, Inches(1.6), Inches(0.85), Inches(10), Inches(0.3),
            "Adapted Scrum methodology — 4 sprints + Sprint 0", font_size=14, color=SUBTLE)

# Accent line under header
add_shape(slide, Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.03), ACCENT)

# ─── Table ──────────────────────────────────────────────────────────────────
sprint_data = [
    ("Sprint 0", "2 weeks", "Stack selection, database schema draft, repository scaffolding, CI setup, environment variables, Neon database provisioned, Pinecone index created."),
    ("Sprint 1", "4 weeks", "Authentication (register/login/me/profile), role hierarchy + JWT, base layout, home page, student and professor dashboards, course CRUD without AI, category seeding."),
    ("Sprint 2", "4 weeks", "Course content builder (section/subsection/lesson blocks), file uploads, enrollment (free), Stripe checkout for paid courses, course feedback, basic notifications."),
    ("Sprint 3", "4 weeks", "AI course generation from PDF (PyMuPDF + Gemini), RAG pipeline (Pinecone + Gemini embeddings), in-course AI tutor, QCM generation, discussions per subsection (posts/replies/votes/reports/AI summary)."),
    ("Sprint 4", "4 weeks", "Gamification engine (XP logs, levels, streaks, achievements, badges), friends + chat (WebSocket), notifications system, learner analytics, professor analytics, multi-tenant admin dashboards, performance pass."),
]

rows = len(sprint_data) + 1
cols = 3
table = add_table(slide, rows, cols, Inches(0.8), Inches(1.7), Inches(11.7), Inches(4.5))

# Column widths
table.columns[0].width = Inches(1.8)
table.columns[1].width = Inches(1.4)
table.columns[2].width = Inches(8.5)

# Header row
header_fill = DARK
headers = ["Sprint", "Duration", "Goals"]
for i, h in enumerate(headers):
    set_cell(table, 0, i, h, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, fill_color=DARK)

# Sprint colors
sprint_colors = [
    RGBColor(0xE8, 0xF0, 0xFE),  # Sprint 0 - light blue
    RGBColor(0xE0, 0xF2, 0xFE),  # Sprint 1
    RGBColor(0xDC, 0xF5, 0xE6),  # Sprint 2 - light green
    RGBColor(0xFE, 0xF3, 0xE2),  # Sprint 3 - light orange
    RGBColor(0xF3, 0xE8, 0xFF),  # Sprint 4 - light purple
]

for idx, (sprint, duration, goals) in enumerate(sprint_data):
    row = idx + 1
    bg = sprint_colors[idx]
    set_cell(table, row, 0, sprint, font_size=11, color=DARK, bold=True, alignment=PP_ALIGN.CENTER, fill_color=bg)
    set_cell(table, row, 1, duration, font_size=11, color=DARK, alignment=PP_ALIGN.CENTER, fill_color=bg)
    set_cell(table, row, 2, goals, font_size=10, color=SUBTLE, fill_color=bg)

# ─── Bottom note ─────────────────────────────────────────────────────────────
add_textbox(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.4),
            "Total duration: 18 weeks | Each sprint delivers a working increment of the platform",
            font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)

# Agile methodology card
add_rounded_rect(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5), CARD_BG, BORDER)
add_textbox(slide, Inches(1.0), Inches(6.85), Inches(11.3), Inches(0.4),
            "Methodology: Adapted Scrum — fixed-length sprints, backlog refinement, sprint review with supervisor, written retrospective per cycle",
            font_size=10, color=SUBTLE, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 1, 1)

# ─── Save ──────────────────────────────────────────────────────────────────
output_path = "Sprint_Plan.pptx"
prs.save(output_path)
print(f"[OK] Sprint plan saved: {output_path}")
