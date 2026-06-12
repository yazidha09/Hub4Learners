from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = RGBColor(0xF1, 0xF5, 0xF9)
DARK    = RGBColor(0x1E, 0x29, 0x3B)
SUBTLE  = RGBColor(0x47, 0x52, 0x66)
BORDER  = RGBColor(0xE2, 0xE8, 0xF0)
ACCENT  = RGBColor(0x25, 0x63, 0xEB)
ACCENT2 = RGBColor(0xF9, 0x73, 0x16)
GREEN   = RGBColor(0x05, 0x9C, 0x69)

logo = "frontend/src/assets/logo/icon.PNG"

def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    return s

def rrect(slide, l, t, w, h, fc, bc=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if bc: s.line.color.rgb = bc; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, sz=18, c=DARK, b=False, al=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
    p.font.name = 'Calibri'; p.alignment = al
    return tb

def header(slide, num, title, accent=ACCENT):
    rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), accent)
    rect(slide, Inches(0.8), Inches(0.4), Inches(0.06), Inches(0.6), accent)
    txt(slide, Inches(1.1), Inches(0.4), Inches(11), Inches(0.5), title, 30, DARK, True)

# ═══════════════════════════════════════════
# SLIDE 1 — Frontend
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide, WHITE)
header(slide, 1, "Frontend Technologies", ACCENT)
if os.path.exists(logo): slide.shapes.add_picture(logo, Inches(12.0), Inches(0.3), Inches(0.7), Inches(0.7))

items = [
    ("React 19 + TypeScript", "UI framework with lazy loading & hooks"),
    ("Vite 7", "Build tool with HMR & code-splitting"),
    ("Tailwind CSS 3", "Utility-first styling with dark mode"),
    ("react-router-dom 7", "SPA routing with lazy pages"),
    ("TipTap Editor", "Rich text (color, align, highlight)"),
    ("react-markdown + react-pdf", "Content rendering"),
]
x1, y1 = Inches(0.8), Inches(1.3)
for i, (title, desc) in enumerate(items):
    col, row = i % 3, i // 3
    x = x1 + col * Inches(4.1); y = y1 + row * Inches(2.8)
    rrect(slide, x, y, Inches(3.7), Inches(2.4), CARD_BG, BORDER)
    rect(slide, x, y, Inches(3.7), Inches(0.06), ACCENT)
    txt(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.1), Inches(0.3), title, 14, DARK, True)
    txt(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.1), Inches(0.8), desc, 11, SUBTLE)

# Contexte/Hooks section at bottom
rrect(slide, Inches(0.8), Inches(6.5), Inches(5.5), Inches(0.7), CARD_BG, BORDER)
rect(slide, Inches(0.8), Inches(6.5), Inches(0.06), Inches(0.7), GREEN)
txt(slide, Inches(1.2), Inches(6.55), Inches(4.8), Inches(0.6),
    "Context & Hooks: AuthContext, GamificationContext, useNotifications, useSettings", 10, SUBTLE)

# API Client bottom
rrect(slide, Inches(6.8), Inches(6.5), Inches(5.7), Inches(0.7), CARD_BG, BORDER)
rect(slide, Inches(6.8), Inches(6.5), Inches(0.06), Inches(0.7), ACCENT2)
txt(slide, Inches(7.2), Inches(6.55), Inches(5.0), Inches(0.6),
    "API Client: auth, course, admin, org, friends, gamification, qcm, discussions, payment, billing, notifications, category, public", 10, SUBTLE)

# ═══════════════════════════════════════════
# SLIDE 2 — Backend
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide, WHITE)
header(slide, 2, "Backend Technologies", ACCENT2)
if os.path.exists(logo): slide.shapes.add_picture(logo, Inches(12.0), Inches(0.3), Inches(0.7), Inches(0.7))

backend_items = [
    ("Core Framework", ["FastAPI + Python 3.11", "Uvicorn ASGI server", "JWT (HS256) + bcrypt", "WebSockets"]),
    ("AI & RAG", ["Google Gemini 3.1 Flash Lite", "Pinecone Vector DB", "PyMuPDF (PDF parsing)", "Gemini Embeddings (768-dim)"]),
    ("Routes (16)", ["Auth, Course, Category, Admin", "Org, AI, Friend, Notifications", "Course-Gen, Announcement", "Payment, Billing, Gamification", "Discussion, Public, WebSocket"]),
    ("Middleware", ["CORS (localhost:5173)", "GZip compression", "Static Files cache", "Connection Pool (20/30)"]),
]
x1, y1 = Inches(0.8), Inches(1.3)
colors = [ACCENT, ACCENT2, GREEN, ACCENT]
for i, (title, bullets) in enumerate(backend_items):
    col = i % 2; row = i // 2
    x = x1 + col * Inches(6.2); y = y1 + row * Inches(2.8)
    c = colors[i]
    rrect(slide, x, y, Inches(5.8), Inches(2.4), CARD_BG, BORDER)
    rect(slide, x, y, Inches(5.8), Inches(0.06), c)
    txt(slide, x + Inches(0.3), y + Inches(0.3), Inches(5.2), Inches(0.3), title, 14, DARK, True)
    by = y + Inches(0.7)
    for b in bullets:
        txt(slide, x + Inches(0.5), by, Inches(5.0), Inches(0.22), f"- {b}", 10, SUBTLE)
        by += Inches(0.28)

# ═══════════════════════════════════════════
# SLIDE 3 — Database & Infrastructure
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide, WHITE)
header(slide, 3, "Database & Infrastructure", GREEN)
if os.path.exists(logo): slide.shapes.add_picture(logo, Inches(12.0), Inches(0.3), Inches(0.7), Inches(0.7))

# Left: Database
rrect(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5), CARD_BG, BORDER)
rect(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.06), GREEN)
txt(slide, Inches(1.1), Inches(1.5), Inches(5.0), Inches(0.3), "Database", 15, DARK, True)

db_items = [
    "PostgreSQL (Neon.tech Serverless)",
    "SQLModel 0.0.16 + SQLAlchemy",
    "asyncpg + psycopg2-binary",
    "25+ tables: users, courses, enrollments,",
    "  gamification, discussions, friendships,",
    "  notifications, payments, analytics",
    "Auto migrations on startup",
    "Vector DB: Pinecone (hub4learners)",
]
dy = Inches(2.0)
for item in db_items:
    txt(slide, Inches(1.3), dy, Inches(4.8), Inches(0.22), f"- {item}", 10, SUBTLE)
    dy += Inches(0.3)

# Right: Infrastructure
rrect(slide, Inches(6.8), Inches(1.3), Inches(5.7), Inches(5.5), CARD_BG, BORDER)
rect(slide, Inches(6.8), Inches(1.3), Inches(5.7), Inches(0.06), ACCENT)
txt(slide, Inches(7.1), Inches(1.5), Inches(5.2), Inches(0.3), "Infrastructure & Services", 15, DARK, True)

infra_items = [
    "Stripe - Checkout & Pro ($9.99/mo)",
    "Google Gemini - LLM + Embeddings",
    "Neon.tech - Serverless PostgreSQL",
    "Git - Version control",
    "UV - Python dependency management",
    "3-Tier: React -> FastAPI -> PostgreSQL",
    "Real-time: WebSocket (messaging, notifs)",
    "Frontend: Vite build, code-splitting",
]
dy = Inches(2.0)
for item in infra_items:
    txt(slide, Inches(7.3), dy, Inches(5.0), Inches(0.22), f"- {item}", 10, SUBTLE)
    dy += Inches(0.3)

output = "Techno_Slides.pptx"
prs.save(output)
print(f"[OK] {output} ({len(prs.slides)} slides)")
