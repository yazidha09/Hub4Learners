from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG    = RGBColor(0xF1, 0xF5, 0xF9)
DARK       = RGBColor(0x1E, 0x29, 0x3B)
SUBTLE     = RGBColor(0x47, 0x52, 0x66)
MUTED      = RGBColor(0x94, 0xA3, 0xB8)
BORDER     = RGBColor(0xE2, 0xE8, 0xF0)
ACCENT     = RGBColor(0x25, 0x63, 0xEB)
ACCENT2    = RGBColor(0xF9, 0x73, 0x16)
GREEN      = RGBColor(0x05, 0x9C, 0x69)
PURPLE     = RGBColor(0x7C, 0x3A, 0xED)
TEAL       = RGBColor(0x0D, 0x94, 0x8A)

ACTOR_COLORS = [MUTED, ACCENT, ACCENT2, GREEN, PURPLE]

logo_path = "frontend/src/assets/logo/icon.PNG"


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(tf, text, font_size=16, color=DARK, bold=False, space_before=Pt(1), space_after=Pt(0), font_name='Calibri'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    return p


actors = [
    {
        "role": "Visitor",
        "tag": "Unauthenticated",
        "icon": "👤",
        "items": [
            "Browse home page & course catalog",
            "Register a new account",
            "Log in to the platform",
        ],
    },
    {
        "role": "Student",
        "tag": "Learner",
        "icon": "🎓",
        "items": [
            "Enroll in courses (free & paid)",
            "Complete lessons & take QCM quizzes",
            "Earn XP, level up, unlock badges",
            "Ask the AI tutor questions",
            "Post in discussions & chat with friends",
        ],
    },
    {
        "role": "Professor",
        "tag": "Content Author",
        "icon": "👨‍🏫",
        "items": [
            "Create, publish & manage courses",
            "Design section / lesson block hierarchy",
            "Upload PDFs for AI course generation",
            "Monitor enrollments & course feedback",
            "Manage chat requests from students",
        ],
    },
    {
        "role": "University Admin",
        "tag": "Institutional Moderator",
        "icon": "🏛️",
        "items": [
            "Review & approve professor join requests",
            "Post university announcements",
            "Moderate discussions & content",
            "Manage users within their university",
        ],
    },
    {
        "role": "Super Admin",
        "tag": "Platform Owner",
        "icon": "⚙️",
        "items": [
            "Create regions & universities",
            "Create University Admin accounts",
            "Ban / unban any user globally",
            "Oversee platform-wide moderation",
        ],
    },
]

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Actors
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

if os.path.exists(logo_path):
    slide.shapes.add_picture(logo_path, Inches(0.8), Inches(0.2), Inches(0.45), Inches(0.45))

add_textbox(slide, Inches(1.4), Inches(0.2), Inches(8), Inches(0.4),
            "Hub4Learners — Actors & Permissions", font_size=22, color=DARK, bold=True)
add_textbox(slide, Inches(1.4), Inches(0.6), Inches(8), Inches(0.25),
            "5 actors  |  Role-based access control (JWT)  |  Multi-tenant hierarchy",
            font_size=11, color=SUBTLE)

add_shape(slide, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.02), BORDER)

# ─── 5 cards in a row ──────────────────────────────────────────────────────
card_w = Inches(2.2)
card_h = Inches(5.2)
gap = Inches(0.2)
total_w = 5 * card_w + 4 * gap
start_x = (Inches(13.333) - total_w) / 2
start_y = Inches(1.3)

for idx, actor in enumerate(actors):
    x = start_x + idx * (card_w + gap)
    y = start_y
    color = ACTOR_COLORS[idx]

    add_rounded_rect(slide, x, y, card_w, card_h, CARD_BG, BORDER)
    add_shape(slide, x, y, card_w, Inches(0.04), color)

    # Icon circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.6), y + Inches(0.2), Inches(0.9), Inches(0.9))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    add_textbox(slide, x + Inches(0.6), y + Inches(0.25), Inches(0.9), Inches(0.8),
                actor["icon"], font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Role name
    add_textbox(slide, x + Inches(0.1), y + Inches(1.25), Inches(2.0), Inches(0.35),
                actor["role"], font_size=15, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

    # Tag
    add_rounded_rect(slide, x + Inches(0.3), y + Inches(1.6), Inches(1.6), Inches(0.25), color)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.62), Inches(1.6), Inches(0.2),
                actor["tag"], font_size=8, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Items
    txBox = add_textbox(slide, x + Inches(0.15), y + Inches(2.05), Inches(1.9), Inches(3.0),
                        "", font_size=10, color=DARK)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = actor["items"][0]
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = DARK
    tf.paragraphs[0].font.name = 'Calibri'
    for item in actor["items"][1:]:
        add_paragraph(tf, item, font_size=9, color=DARK)

# Bottom note
add_rounded_rect(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.45), CARD_BG, BORDER)
add_textbox(slide, Inches(1.0), Inches(6.85), Inches(11.3), Inches(0.35),
            "Access control enforced via JWT role claims + FastAPI middleware  |  Hierarchy: Super Admin > University Admin > Professor > Student > Visitor",
            font_size=10, color=MUTED, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
            "1/1", font_size=10, color=MUTED, alignment=PP_ALIGN.RIGHT)

output_path = "Actors_Overview.pptx"
prs.save(output_path)
print(f"[OK] Saved: {output_path}")
