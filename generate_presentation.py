from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Color palette (white background) ────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT   = RGBColor(0xF8, 0xFA, 0xFC)
CARD_BG    = RGBColor(0xF1, 0xF5, 0xF9)
DARK       = RGBColor(0x1E, 0x29, 0x3B)
SUBTLE     = RGBColor(0x47, 0x52, 0x66)
MUTED      = RGBColor(0x94, 0xA3, 0xB8)
BORDER     = RGBColor(0xE2, 0xE8, 0xF0)
ACCENT     = RGBColor(0x25, 0x63, 0xEB)  # blue
ACCENT2    = RGBColor(0xF9, 0x73, 0x16)  # orange
GREEN      = RGBColor(0x05, 0x9C, 0x69)
YELLOW     = RGBColor(0xD9, 0x77, 0x06)

logo_path = "frontend/src/assets/logo/icon.PNG"
logo_horizontal = "frontend/src/assets/logo/Horizontal.PNG"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(tf, text, font_size=16, color=DARK, bold=False, space_before=Pt(4), space_after=Pt(2), font_name='Calibri', alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    p.alignment = alignment
    return p


def add_page_number(slide, num, total=13):
    add_textbox(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                f"{num}/{total}", font_size=10, color=MUTED, alignment=PP_ALIGN.RIGHT)


def section_header(slide, number, title, subtitle=None):
    # Thin accent bar at top
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
    # Section number
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(1.5), Inches(0.4),
                f"0{number}" if number < 10 else str(number), font_size=36, color=ACCENT, bold=True)
    # Title
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
                title, font_size=28, color=DARK, bold=True)
    # Subtitle
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.4),
                    subtitle, font_size=14, color=SUBTLE)


TOTAL_SLIDES = 13

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

# Top accent bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)

# Logo
if os.path.exists(logo_path):
    slide.shapes.add_picture(logo_path, Inches(0.8), Inches(0.5), Inches(0.7), Inches(0.7))

# Title block
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(1.0),
            "Hub4Learners", font_size=56, color=DARK, bold=True)
add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(0.5),
            "Intelligent & Interactive Learning Platform", font_size=22, color=SUBTLE)
add_textbox(slide, Inches(0.8), Inches(3.5), Inches(11), Inches(0.4),
            "AI-powered course generation, gamification, and real-time community", font_size=14, color=MUTED)

# Accent line
add_shape(slide, Inches(0.8), Inches(4.2), Inches(3), Inches(0.04), ACCENT)

# Team info
txBox = add_textbox(slide, Inches(0.8), Inches(4.7), Inches(11), Inches(1.5),
                    "", font_size=14, color=SUBTLE)
tf = txBox.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Final Year Project"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = MUTED
tf.paragraphs[0].font.name = 'Calibri'

add_paragraph(tf, "Team : [Member Names]", font_size=18, color=DARK, bold=True)
add_paragraph(tf, "Supervisor : [Supervisor Name]", font_size=14, color=SUBTLE)
add_paragraph(tf, "Institution : [School/University Name]  |  Academic Year 2025-2026", font_size=14, color=SUBTLE)

add_page_number(slide, 1, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Table of Contents
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
section_header(slide, 2, "Table of Contents")

items = [
    ("01", "Context & Problem Statement", "Why Hub4Learners?"),
    ("02", "Project Objectives", "What we aim to achieve"),
    ("03", "Requirements Analysis", "Use cases and specifications"),
    ("04", "Architecture & Design", "Software architecture and data model"),
    ("05", "Technology Stack", "Complete technical stack"),
    ("06", "Implementation - Backend", "API, AI, RAG, WebSocket"),
    ("07", "Implementation - Frontend", "User interface and components"),
    ("08", "Key Features", "Gamification, Payments, Social"),
    ("09", "Demonstration", "Application walkthrough"),
    ("10", "Testing & Quality", "Project validation"),
    ("11", "Conclusion & Perspectives", "Challenges and future improvements"),
]

y = Inches(2.3)
for num, title, desc in items:
    add_textbox(slide, Inches(1.2), y, Inches(0.4), Inches(0.35),
                num, font_size=16, color=ACCENT, bold=True)
    add_textbox(slide, Inches(1.8), y, Inches(4.5), Inches(0.25),
                title, font_size=14, color=DARK, bold=True)
    add_textbox(slide, Inches(1.8), y + Inches(0.2), Inches(4.5), Inches(0.2),
                desc, font_size=11, color=MUTED)
    y += Inches(0.42)

add_page_number(slide, 2, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Context & Problem Statement
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
section_header(slide, 3, "Context & Problem Statement")

# Card 1
add_rounded_rect(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(2.4), CARD_BG, BORDER)
add_shape(slide, Inches(0.8), Inches(2.4), Inches(0.06), Inches(2.4), ACCENT)
add_textbox(slide, Inches(1.2), Inches(2.6), Inches(5), Inches(0.35),
            "Observation", font_size=18, color=ACCENT, bold=True)
txBox = add_textbox(slide, Inches(1.2), Inches(3.1), Inches(5), Inches(1.5),
                    "", font_size=14, color=DARK)
tf = txBox.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Existing learning platforms often suffer from:"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = DARK
tf.paragraphs[0].font.name = 'Calibri'
add_paragraph(tf, "- Rigidity: static content, little interactivity", font_size=13, color=SUBTLE)
add_paragraph(tf, "- Genericity: no curriculum-adapted content", font_size=13, color=SUBTLE)
add_paragraph(tf, "- Isolation: no community or peer support", font_size=13, color=SUBTLE)
add_paragraph(tf, "- Cost: expensive subscriptions for students", font_size=13, color=SUBTLE)

# Card 2
add_rounded_rect(slide, Inches(6.8), Inches(2.4), Inches(5.5), Inches(2.4), CARD_BG, BORDER)
add_shape(slide, Inches(6.8), Inches(2.4), Inches(0.06), Inches(2.4), ACCENT2)
add_textbox(slide, Inches(7.2), Inches(2.6), Inches(5), Inches(0.35),
            "Problem Statement", font_size=18, color=ACCENT2, bold=True)
txBox = add_textbox(slide, Inches(7.2), Inches(3.1), Inches(5), Inches(1.5),
                    "", font_size=14, color=DARK)
tf = txBox.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "How to design a platform that:"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = DARK
tf.paragraphs[0].font.name = 'Calibri'
add_paragraph(tf, "- Lets professors create courses easily", font_size=13, color=SUBTLE)
add_paragraph(tf, "- Provides AI tutoring based on course content", font_size=13, color=SUBTLE)
add_paragraph(tf, "- Gamifies learning to maintain motivation", font_size=13, color=SUBTLE)
add_paragraph(tf, "- Facilitates social interactions and communication", font_size=13, color=SUBTLE)

add_page_number(slide, 3, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Project Objectives
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
section_header(slide, 4, "Project Objectives")

objectives = [
    ("Main Objective", "Build an interactive LMS platform where professors publish structured content and students learn with AI assistance, progress tracking, and an active community."),
    ("Auth & Profile", "Secure authentication (JWT), roles (student, professor, admin), user profiles with bio and avatar."),
    ("Course Creation", "Rich editor (TipTap), Course - Section - Subsection - Lesson Block hierarchy, media (images, videos, PDFs)."),
    ("AI & RAG", "Contextual AI tutor (Gemini + Pinecone), MCQ generation, course summaries, PDF-to-course generation."),
    ("Gamification", "XP, levels, streaks, achievements, badges, daily/weekly/all-time leaderboards."),
    ("Social & Communication", "Per-lesson discussions, real-time messaging (WebSocket), live notifications, university announcements."),
    ("Payments & Pro", "Stripe Checkout for paid courses, Pro subscription ($9.99/month) for advanced features."),
    ("Administration", "Admin dashboard, user/role management, platform statistics, university management."),
]

y = Inches(2.3)
for i, (title, desc) in enumerate(objectives):
    # Numbered circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y + Pt(2), Inches(0.28), Inches(0.28))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    circ.text_frame.text = str(i + 1)
    circ.text_frame.paragraphs[0].font.size = Pt(10)
    circ.text_frame.paragraphs[0].font.color.rgb = WHITE
    circ.text_frame.paragraphs[0].font.bold = True
    circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Title
    add_textbox(slide, Inches(1.4), y - Pt(2), Inches(2.2), Inches(0.25),
                title, font_size=12, color=ACCENT, bold=True)
    # Description
    add_textbox(slide, Inches(3.6), y - Pt(2), Inches(8.5), Inches(0.35),
                desc, font_size=11, color=SUBTLE)
    y += Inches(0.55)

add_page_number(slide, 4, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Requirements Analysis
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
section_header(slide, 5, "Requirements Analysis", "Use case diagram - Actors and features")

# Actors column
add_rounded_rect(slide, Inches(0.8), Inches(2.3), Inches(2.8), Inches(4.2), CARD_BG, BORDER)
add_textbox(slide, Inches(1.0), Inches(2.5), Inches(2.4), Inches(0.3),
            "Actors", font_size=16, color=ACCENT, bold=True)

actors = [
    ("Visitor", "Register, browse catalog"),
    ("Student", "Learn, QCM, XP, discussions"),
    ("Professor", "Create courses, analytics"),
    ("University Admin", "Manage users/announcements"),
    ("Super Admin", "Manage platform & universities"),
]

y = Inches(3.0)
for actor, desc in actors:
    add_textbox(slide, Inches(1.1), y, Inches(2.3), Inches(0.22),
                actor, font_size=12, color=DARK, bold=True)
    add_textbox(slide, Inches(1.1), y + Inches(0.2), Inches(2.3), Inches(0.2),
                desc, font_size=10, color=MUTED)
    y += Inches(0.5)

# Features column
add_rounded_rect(slide, Inches(4.0), Inches(2.3), Inches(8.3), Inches(4.2), CARD_BG, BORDER)
add_textbox(slide, Inches(4.3), Inches(2.5), Inches(7.7), Inches(0.3),
            "Main Features (by sprint)", font_size=16, color=ACCENT2, bold=True)

usecases = [
    ("Sprint 1 - Auth", "Registration, login, JWT, roles (5 use cases)"),
    ("Sprint 2 - Courses", "CRUD courses, sections, lessons, enrollment, progress (6)"),
    ("Sprint 3 - AI", "RAG tutor, QCM, summaries, PDF to course gen (6)"),
    ("Sprint 4 - Gamification", "XP, levels, streaks, achievements, badges, leaderboards (11)"),
    ("Sprint 5 - Communication", "Discussions, real-time messaging, notifications (9)"),
    ("Sprint 6 - Payments & Admin", "Stripe, analytics, admin dashboard, universities (10)"),
]

y = Inches(3.0)
for sprint, desc in usecases:
    add_textbox(slide, Inches(4.3), y, Inches(3), Inches(0.22),
                sprint, font_size=11, color=GREEN, bold=True)
    add_textbox(slide, Inches(7.5), y, Inches(4.5), Inches(0.22),
                desc, font_size=11, color=SUBTLE)
    y += Inches(0.45)

add_page_number(slide, 5, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Architecture & Design
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
section_header(slide, 6, "Architecture & Design", "3-tier architecture + external services")

# Architecture tiers
tiers = [
    ("Frontend (React 19 + Vite 7)", Inches(0.8), Inches(2.4),
     "Pages, Components, Contexts, Hooks\nTypeScript, Tailwind CSS, TipTap", ACCENT),
    ("Backend (FastAPI + Python 3.11)", Inches(0.8), Inches(3.8),
     "16 routers, Controllers, Utilities\nSQLModel ORM, WebSocket, JWT", ACCENT2),
    ("Database (PostgreSQL + Neon.tech)", Inches(0.8), Inches(5.2),
     "20 models, auto migrations, indexes\nPinecone (vectors), Stripe (payments)", GREEN),
]

for i, (title, left, top, desc, color) in enumerate(tiers):
    add_rounded_rect(slide, left, top, Inches(5.8), Inches(1.2), CARD_BG, BORDER)
    add_shape(slide, left, top, Inches(0.06), Inches(1.2), color)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.15), Inches(5.2), Inches(0.3),
                title, font_size=13, color=DARK, bold=True)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.5), Inches(5.2), Inches(0.55),
                desc, font_size=11, color=SUBTLE)
    if i < len(tiers) - 1:
        add_textbox(slide, left + Inches(2.5), top + Inches(1.15), Inches(2.5), Inches(0.2),
                    "REST API / WebSocket", font_size=9, color=MUTED, alignment=PP_ALIGN.CENTER)

# External services
add_rounded_rect(slide, Inches(7.2), Inches(2.4), Inches(5.3), Inches(2.0), CARD_BG, BORDER)
add_textbox(slide, Inches(7.5), Inches(2.6), Inches(4.8), Inches(0.3),
            "External Services", font_size=16, color=ACCENT, bold=True)

services = [
    ("Google Gemini", "LLM for chat, QCM, summaries, embeddings"),
    ("Pinecone", "Vector database for RAG (768 dimensions)"),
    ("Stripe", "Checkout sessions, payment validation"),
    ("Neon.tech", "Serverless PostgreSQL, pool size 20"),
]

y = Inches(3.1)
for service, desc in services:
    add_textbox(slide, Inches(7.6), y, Inches(2.3), Inches(0.22),
                service, font_size=11, color=DARK, bold=True)
    add_textbox(slide, Inches(10.0), y, Inches(2.3), Inches(0.22),
                desc, font_size=10, color=SUBTLE)
    y += Inches(0.38)

# Data model
add_rounded_rect(slide, Inches(7.2), Inches(4.7), Inches(5.3), Inches(1.7), CARD_BG, BORDER)
add_textbox(slide, Inches(7.5), Inches(4.9), Inches(4.8), Inches(0.3),
            "Data Model (20 entities)", font_size=14, color=ACCENT2, bold=True)
add_textbox(slide, Inches(7.6), Inches(5.3), Inches(4.7), Inches(0.9),
            "User - University | Course - Section - Subsection - LessonBlock\n"
            "Enrollment | CourseProgress | CourseFeedback | QCMAttempt\n"
            "Friendship - FriendMessage | DiscussionPost - Vote/Report/Summary\n"
            "UserGamification - XPLog | Achievement | Badge\n"
            "Notification | Announcement | GeneratedCourse | Category",
            font_size=10, color=SUBTLE)

add_page_number(slide, 6, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Technology Stack
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
section_header(slide, 7, "Technology Stack")

techs = [
    ("Backend", [
        ("FastAPI", "Async Python web framework"),
        ("SQLModel / SQLAlchemy", "ORM and data models"),
        ("PostgreSQL (Neon.tech)", "Serverless database"),
        ("asyncpg / psycopg2", "PostgreSQL connectors"),
        ("python-jose / bcrypt", "JWT and password hashing"),
        ("Google Generative AI", "Gemini API, embeddings"),
        ("Pinecone", "Vector database for RAG"),
        ("PyMuPDF (fitz)", "PDF parsing into chunks"),
        ("Stripe SDK", "Payments and checkout sessions"),
    ]),
    ("Frontend", [
        ("React 19 + TypeScript", "Typed UI framework"),
        ("Vite 7", "Fast build tool"),
        ("Tailwind CSS 3", "Utility-first styles"),
        ("react-router-dom 7", "SPA routing"),
        ("TipTap Editor", "Rich editor (color, align, highlight)"),
        ("react-markdown", "Markdown rendering"),
        ("react-pdf", "PDF display"),
        ("WebSocket API", "Real-time messaging"),
        ("Code-splitting", "React.lazy for optimized bundles"),
    ]),
    ("Services & Tools", [
        ("Neon.tech", "Managed serverless PostgreSQL"),
        ("Pinecone", "Vector index (hub4learners)"),
        ("Google Gemini", "LLM + embeddings (768-dim)"),
        ("Stripe", "Payment checkout"),
        ("Git", "Version control"),
        ("UV (Python)", "Dependency management"),
    ]),
]

x_positions = [Inches(0.8), Inches(5.0), Inches(9.2)]
for col_idx, (col_title, items) in enumerate(techs):
    add_rounded_rect(slide, x_positions[col_idx], Inches(2.3), Inches(3.8), Inches(4.6), CARD_BG, BORDER)
    add_shape(slide, x_positions[col_idx], Inches(2.3), Inches(3.8), Inches(0.06), ACCENT)
    add_textbox(slide, x_positions[col_idx] + Inches(0.3), Inches(2.5), Inches(3.2), Inches(0.3),
                col_title, font_size=15, color=DARK, bold=True)

    y = Inches(3.0)
    for tech, desc in items:
        add_textbox(slide, x_positions[col_idx] + Inches(0.3), y, Inches(3.2), Inches(0.22),
                    tech, font_size=10, color=DARK, bold=True)
        add_textbox(slide, x_positions[col_idx] + Inches(0.3), y + Inches(0.18), Inches(3.2), Inches(0.2),
                    desc, font_size=9, color=SUBTLE)
        y += Inches(0.4)

add_page_number(slide, 7, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Implementation Backend
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
section_header(slide, 8, "Implementation - Backend", "16 API routers, 4 AI utilities, real-time WebSocket")

features_backend = [
    ("Auth & Security", "JWT HS256, bcrypt, roles (student to super_admin), authorization middleware"),
    ("Course Management", "Full CRUD, 4-level hierarchy, media upload, progress tracking, feedback"),
    ("AI & RAG", "Contextual tutor (Gemini to Pinecone), MCQ generation, course summaries, PDF to course"),
    ("Gamification", "Centralized XP service (anti-cheat, daily caps), achievements, badges, leaderboards"),
    ("Communication", "Per-lesson discussions (vote, report), AI thread summaries"),
    ("Social", "User search, friend requests, real-time text/media messaging"),
    ("Notifications", "Push + WebSocket notifications, read/delete, university announcements"),
    ("Organizations", "University CRUD, professor join requests, user reassignment"),
    ("Payments", "Stripe Checkout, frontend confirmation to enrollment, Pro at $9.99/month"),
    ("Analytics", "Professor analytics (at-risk learners), student analytics (progress, QCM, activity)"),
]

y = Inches(2.3)
for i, (title, desc) in enumerate(features_backend):
    row_color = ACCENT if i % 2 == 0 else ACCENT2
    add_shape(slide, Inches(0.9), y, Inches(0.04), Inches(0.35), row_color)
    add_textbox(slide, Inches(1.2), y, Inches(2.5), Inches(0.35),
                title, font_size=11, color=DARK, bold=True)
    add_textbox(slide, Inches(3.8), y, Inches(8.5), Inches(0.35),
                desc, font_size=11, color=SUBTLE)
    y += Inches(0.45)

add_page_number(slide, 8, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Implementation Frontend
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
section_header(slide, 9, "Implementation - Frontend", "9 pages, 20 components, contexts, hooks, and API client")

# Pages section
add_rounded_rect(slide, Inches(0.8), Inches(2.3), Inches(5.8), Inches(2.0), CARD_BG, BORDER)
add_textbox(slide, Inches(1.1), Inches(2.5), Inches(5), Inches(0.3),
            "Pages (lazy-loaded)", font_size=14, color=ACCENT, bold=True)

pages_list = [
    "HomePage - Landing with live stats, categories, CTA",
    "LoginPage / RegisterPage - Auth with validation",
    "DashboardPage - Central hub, role-based redirect",
    "StudentDashboard - Enrolled courses, progress, stats",
    "ProfessorDashboard - Course management, analytics, PDF gen",
    "AdminDashboard - Stats, users, courses, universities",
    "CourseLearningPage - Course reader, AI tutor, QCM, discussions",
    "PaymentResultPage - Stripe success/cancel confirmation",
]

y = Inches(2.9)
for p in pages_list:
    add_textbox(slide, Inches(1.3), y, Inches(5), Inches(0.2),
                p, font_size=10, color=SUBTLE)
    y += Inches(0.25)

# Components section
add_rounded_rect(slide, Inches(7.0), Inches(2.3), Inches(5.5), Inches(2.7), CARD_BG, BORDER)
add_textbox(slide, Inches(7.3), Inches(2.5), Inches(5), Inches(0.3),
            "Key Components", font_size=14, color=ACCENT2, bold=True)

components_list = [
    "DashboardLayout - Navigation, theme, notifications, profile",
    "DiscussionSection - Posts, votes, reports, AI summary",
    "QCMModal - Difficulty selection, quiz, result, XP",
    "FriendChat / FriendsMessenger - Real-time WebSocket chat",
    "RichTextEditor - TipTap (color, align, highlight)",
    "UpgradeProModal - Stripe checkout, feature list",
    "GamificationToasts - XP, level, achievement animations",
    "Leaderboard / XPBar / StreakWidget / BadgeShowcase",
]

y = Inches(2.9)
for c in components_list:
    add_textbox(slide, Inches(7.4), y, Inches(4.8), Inches(0.2),
                c, font_size=10, color=SUBTLE)
    y += Inches(0.28)

# Contexts & Hooks
add_rounded_rect(slide, Inches(0.8), Inches(4.6), Inches(5.8), Inches(0.9), CARD_BG, BORDER)
add_textbox(slide, Inches(1.1), Inches(4.7), Inches(5.3), Inches(0.25),
            "Contexts & Hooks", font_size=12, color=GREEN, bold=True)
add_textbox(slide, Inches(1.1), Inches(5.0), Inches(5.3), Inches(0.4),
            "AuthContext (token, role)  |  GamificationContext (XP toast queue)\nuseNotifications (REST + WebSocket)  |  useSettings (theme, compact, privacy)",
            font_size=10, color=SUBTLE)

# API Client
add_rounded_rect(slide, Inches(7.0), Inches(5.2), Inches(5.5), Inches(0.9), CARD_BG, BORDER)
add_textbox(slide, Inches(7.3), Inches(5.3), Inches(5), Inches(0.25),
            "API Client", font_size=12, color=GREEN, bold=True)
add_textbox(slide, Inches(7.3), Inches(5.6), Inches(5), Inches(0.4),
            "_client.ts - dedupGet, cachedGet, invalidate\nPer-resource modules: auth, course, admin, qcm, friends, gamification, billing, payment, org, public, notifications, discussions, category",
            font_size=10, color=SUBTLE)

add_page_number(slide, 9, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Demonstration
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
section_header(slide, 10, "Demonstration")

flow_steps = [
    ("1", "Sign-up / Login", "Create student or professor account"),
    ("2", "Course Discovery", "Browse catalog, filter by category"),
    ("3", "Learning", "Course navigation, lessons, progress tracking"),
    ("4", "AI Tutor", "Ask a question, contextualized answer (RAG)"),
    ("5", "Quiz & Assessment", "MCQ generation, submission, score + XP"),
    ("6", "Gamification", "XP, levels, streaks, badges, leaderboards"),
    ("7", "Discussion & Help", "Posts, votes, AI summary per lesson"),
    ("8", "Messaging", "Real-time chat between friends (WebSocket)"),
    ("9", "Course Creation (Prof)", "Rich editor, PDF upload, AI generation"),
    ("10", "Payment & Pro (if applicable)", "Stripe Checkout, Pro subscription"),
]

y = Inches(2.3)
for i, (num, title, desc) in enumerate(flow_steps):
    col = 0 if i < 5 else 1
    x_base = Inches(0.8) if col == 0 else Inches(6.8)
    row = i % 5
    y_pos = y + row * Inches(0.85)

    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x_base, y_pos + Pt(2), Inches(0.32), Inches(0.32))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    circ.text_frame.text = num
    circ.text_frame.paragraphs[0].font.size = Pt(11)
    circ.text_frame.paragraphs[0].font.color.rgb = WHITE
    circ.text_frame.paragraphs[0].font.bold = True
    circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, x_base + Inches(0.5), y_pos, Inches(2.5), Inches(0.22),
                title, font_size=12, color=DARK, bold=True)
    add_textbox(slide, x_base + Inches(0.5), y_pos + Inches(0.22), Inches(5), Inches(0.2),
                desc, font_size=10, color=SUBTLE)

add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.3),
            "Launch app at localhost:5173  |  API at localhost:8000/api",
            font_size=13, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 10, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Testing, Quality & Deployment
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
section_header(slide, 11, "Testing, Quality & Deployment")

# Left column
add_rounded_rect(slide, Inches(0.8), Inches(2.3), Inches(6.0), Inches(4.5), CARD_BG, BORDER)
add_textbox(slide, Inches(1.1), Inches(2.5), Inches(5.4), Inches(0.3),
            "Software Quality", font_size=16, color=ACCENT, bold=True)

quality_items = [
    ("Security", [
        "Password hashing (bcrypt)",
        "JWT with 60-minute expiry",
        "Role-based authorization middleware",
        "Anti-XSS protection (input validation)",
        "Profanity filter moderation",
    ]),
    ("Performance", [
        "Frontend code-splitting (React.lazy)",
        "GZip compression middleware",
        "Immutable cache for uploads",
        "GET request deduplication (dedupGet)",
        "TTL cache for AI responses",
    ]),
    ("Reliability", [
        "Pinecone lazy-init with lock",
        "Auto RAG re-index (self-heal)",
        "XP anti-cheat (daily cap, one-shot, cooldown)",
        "Auto migrations on startup",
        "Configuration logs at boot",
    ]),
]

y = Inches(2.9)
for cat_title, items in quality_items:
    add_textbox(slide, Inches(1.3), y, Inches(5), Inches(0.22),
                cat_title, font_size=12, color=ACCENT2, bold=True)
    y += Inches(0.25)
    for item in items:
        add_textbox(slide, Inches(1.5), y, Inches(5), Inches(0.18),
                    f"- {item}", font_size=10, color=SUBTLE)
        y += Inches(0.2)
    y += Inches(0.1)

# Right column
add_rounded_rect(slide, Inches(7.2), Inches(2.3), Inches(5.3), Inches(2.0), CARD_BG, BORDER)
add_textbox(slide, Inches(7.5), Inches(2.5), Inches(4.7), Inches(0.3),
            "Deployment", font_size=16, color=GREEN, bold=True)

deploy_items = [
    "Frontend: Vite build to static files",
    "Backend: FastAPI + Uvicorn/Gunicorn",
    "Database: Neon.tech PostgreSQL",
    "API Keys: Gemini, Pinecone, Stripe (.env)",
    "Environment: Windows (dev), Linux (prod)",
]

y = Inches(2.9)
for item in deploy_items:
    add_textbox(slide, Inches(7.6), y, Inches(4.5), Inches(0.22),
                f"- {item}", font_size=11, color=SUBTLE)
    y += Inches(0.3)

# Key Metrics
add_rounded_rect(slide, Inches(7.2), Inches(4.6), Inches(5.3), Inches(2.2), CARD_BG, BORDER)
add_textbox(slide, Inches(7.5), Inches(4.8), Inches(4.7), Inches(0.3),
            "Key Metrics", font_size=16, color=YELLOW, bold=True)

metrics = [
    ("~10,000+", "Lines of code (Python + TypeScript)"),
    ("20", "Database models"),
    ("16", "API routers"),
    ("9", "Frontend pages"),
    ("20+", "React components"),
    ("6", "Development sprints"),
    ("2", "AI services (Gemini + Pinecone)"),
]

y = Inches(5.2)
for val, label in metrics:
    add_textbox(slide, Inches(7.6), y, Inches(0.8), Inches(0.22),
                val, font_size=11, color=ACCENT, bold=True)
    add_textbox(slide, Inches(8.5), y, Inches(3.7), Inches(0.22),
                label, font_size=10, color=SUBTLE)
    y += Inches(0.22)

add_page_number(slide, 11, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Conclusion & Perspectives
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
section_header(slide, 12, "Conclusion & Perspectives")

# Achieved card
add_rounded_rect(slide, Inches(0.8), Inches(2.3), Inches(5.8), Inches(2.8), CARD_BG, BORDER)
add_shape(slide, Inches(0.8), Inches(2.3), Inches(0.06), Inches(2.8), ACCENT)
add_textbox(slide, Inches(1.2), Inches(2.5), Inches(5.2), Inches(0.3),
            "Achieved Objectives", font_size=16, color=ACCENT, bold=True)

bilan_items = [
    "Complete LMS platform delivered in 6 sprints",
    "Integrated generative AI (RAG tutor, QCM, summaries, PDF to course)",
    "Full gamification system (XP, streaks, badges, leaderboards)",
    "Real-time communication (WebSocket, discussions, notifications)",
    "Stripe payments and Pro subscription functional",
    "Modular and extensible architecture",
]

y = Inches(2.9)
for item in bilan_items:
    add_textbox(slide, Inches(1.3), y, Inches(5.2), Inches(0.25),
                f"- {item}", font_size=11, color=SUBTLE)
    y += Inches(0.3)

# Challenges card
add_rounded_rect(slide, Inches(7.0), Inches(2.3), Inches(5.5), Inches(2.8), CARD_BG, BORDER)
add_shape(slide, Inches(7.0), Inches(2.3), Inches(0.06), Inches(2.8), ACCENT2)
add_textbox(slide, Inches(7.4), Inches(2.5), Inches(4.9), Inches(0.3),
            "Challenges Encountered", font_size=16, color=ACCENT2, bold=True)

diff_items = [
    "Pinecone + Gemini integration (768 dimensions, 0.40 score threshold)",
    "Faithful PDF parsing without content loss (2,000 char chunks)",
    "WebSocket notifications + chat with auto-reconnect",
    "XP balancing (daily cap, anti-cheat, cooldowns)",
    "PostgreSQL schema migration (enum to varchar, legacy columns)",
]

y = Inches(2.9)
for item in diff_items:
    add_textbox(slide, Inches(7.4), y, Inches(4.9), Inches(0.25),
                f"- {item}", font_size=11, color=SUBTLE)
    y += Inches(0.3)

# Future Improvements
add_shape(slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.04), BORDER)

add_textbox(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.3),
            "Future Improvements", font_size=16, color=ACCENT, bold=True)

perspectives = [
    "JWT refresh tokens and persistent sessions",
    "Stripe webhook for refunds and disputes",
    "Drag-and-drop course/section reorganization",
    "Mobile push notifications (offline)",
    "Externalized achievement/badge rules (config to DB)",
    "Docker deployment + automated CI/CD",
]

x = Inches(1.0)
y = Inches(6.0)
for i, item in enumerate(perspectives):
    add_textbox(slide, x, y, Inches(3.8), Inches(0.22),
                f"- {item}", font_size=11, color=SUBTLE)
    x += Inches(3.9)
    if (i + 1) % 3 == 0:
        x = Inches(1.0)
        y += Inches(0.3)

add_page_number(slide, 12, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Q&A
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)

if os.path.exists(logo_path):
    slide.shapes.add_picture(logo_path, Inches(5.8), Inches(1.5), Inches(1.5), Inches(1.5))

add_textbox(slide, Inches(0.8), Inches(3.3), Inches(11.5), Inches(0.8),
            "Thank You for Your Attention", font_size=40, color=DARK, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(4.1), Inches(11.5), Inches(0.5),
            "Questions?", font_size=24, color=ACCENT, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.4),
            "Hub4Learners - Intelligent & Interactive Learning Platform",
            font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

# Contact info
if os.path.exists(logo_horizontal):
    slide.shapes.add_picture(logo_horizontal, Inches(4.5), Inches(5.5), Inches(4), Inches(0.5))

add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.3),
            "[Member Names]  |  [Email]  |  [Institution]",
            font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 13, TOTAL_SLIDES)

# ─── Save ──────────────────────────────────────────────────────────────────
output_path = "Hub4Learners_Presentation.pptx"
prs.save(output_path)
print(f"[OK] Presentation saved: {output_path}")
print(f"     Total slides: {len(prs.slides)}")
