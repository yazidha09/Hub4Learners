from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG    = RGBColor(0xF1, 0xF5, 0xF9)
DARK       = RGBColor(0x1E, 0x29, 0x3B)
SUBTLE     = RGBColor(0x47, 0x52, 0x66)
MUTED      = RGBColor(0x94, 0xA3, 0xB8)
BORDER     = RGBColor(0xE2, 0xE8, 0xF0)
ACCENT     = RGBColor(0x25, 0x63, 0xEB)
ACCENT2    = RGBColor(0xF9, 0x73, 0x16)
GREEN      = RGBColor(0x05, 0x9C, 0x69)
PURPLE     = RGBColor(0x7C, 0x3A, 0xED)
TEAL       = RGBColor(0x0D, 0x94, 0x8A)
RED        = RGBColor(0xDC, 0x26, 0x26)

SPRINT_COLORS = [ACCENT, ACCENT2, GREEN, PURPLE, TEAL, RED]

logo_path = "frontend/src/assets/logo/icon.PNG"


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(tf, text, font_size=16, color=DARK, bold=False, space_before=Pt(1), space_after=Pt(0), font_name='Calibri'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    return p


sprints = [
    {
        "id": "Sprint 1",
        "title": "Identity Foundation",
        "duration": "2 weeks",
        "items": [
            "JWT auth, roles & profiles",
            "Base layout & navigation",
            "Home page & course catalog",
            "Student & Professor dashboards",
        ],
    },
    {
        "id": "Sprint 2",
        "title": "Course Architecture",
        "duration": "2 weeks",
        "items": [
            "Course CRUD & categories",
            "Section / Lesson block hierarchy",
            "Rich editor (TipTap)",
            "File uploads & enrollment",
        ],
    },
    {
        "id": "Sprint 3",
        "title": "AI Learning Features",
        "duration": "2 weeks",
        "items": [
            "PDF-to-course generation",
            "RAG pipeline (Pinecone + Gemini)",
            "In-course AI tutor",
            "QCM generation & discussions",
        ],
    },
    {
        "id": "Sprint 4",
        "title": "Gamification & Socials",
        "duration": "2 weeks",
        "items": [
            "XP logging & leveling curve",
            "Daily streaks & achievements",
            "Badge system (event-driven)",
            "Leaderboards & friends",
        ],
    },
    {
        "id": "Sprint 5",
        "title": "Communication & Community",
        "duration": "2 weeks",
        "items": [
            "Real-time chat (WebSocket)",
            "Discussion forums with votes",
            "Content reporting",
            "Live notifications",
        ],
    },
    {
        "id": "Sprint 6",
        "title": "Payment, Analytics & Admin",
        "duration": "2 weeks",
        "items": [
            "Stripe Checkout & webhooks",
            "Pro subscription ($9.99/mo)",
            "Learner & professor analytics",
            "Multi-tenant admin dashboards",
        ],
    },
]

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Sprint Overview (all 6)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = WHITE

# Top accent bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

# Logo
if os.path.exists(logo_path):
    slide.shapes.add_picture(logo_path, Inches(0.8), Inches(0.2), Inches(0.45), Inches(0.45))

# Title
add_textbox(slide, Inches(1.4), Inches(0.2), Inches(8), Inches(0.4),
            "Hub4Learners — Sprint Overview", font_size=22, color=DARK, bold=True)
add_textbox(slide, Inches(1.4), Inches(0.6), Inches(8), Inches(0.25),
            "Adapted Scrum  |  6 sprints × 2 weeks  |  12 weeks total", font_size=11, color=SUBTLE)

# Separator
add_shape(slide, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.02), BORDER)

# ─── 2 rows × 3 columns grid of sprint cards ────────────────────────────
card_w = Inches(3.7)
card_h = Inches(2.85)
gap_x = Inches(0.3)
gap_y = Inches(0.25)
start_x = Inches(0.8)
start_y = Inches(1.3)

positions = []
for row in range(2):
    for col in range(3):
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        positions.append((x, y))

for idx, sprint in enumerate(sprints):
    x, y = positions[idx]
    color = SPRINT_COLORS[idx]

    # Card background
    add_rounded_rect(slide, x, y, card_w, card_h, CARD_BG, BORDER)
    # Left accent stripe
    add_shape(slide, x, y, Inches(0.05), card_h, color)

    # Sprint ID badge
    badge = add_rounded_rect(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.95), Inches(0.3), color)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.17), Inches(0.95), Inches(0.25),
                sprint["id"], font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Title
    add_textbox(slide, x + Inches(1.3), y + Inches(0.12), Inches(2.2), Inches(0.3),
                sprint["title"], font_size=14, color=DARK, bold=True)

    # Duration
    add_textbox(slide, x + Inches(1.3), y + Inches(0.42), Inches(2.2), Inches(0.2),
                sprint["duration"], font_size=9, color=color, bold=True)

    # Items
    txBox = add_textbox(slide, x + Inches(0.25), y + Inches(0.6), Inches(3.2), Inches(2.1),
                        "", font_size=10, color=DARK)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = sprint["items"][0]
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = DARK
    tf.paragraphs[0].font.name = 'Calibri'
    for item in sprint["items"][1:]:
        add_paragraph(tf, item, font_size=10, color=DARK)

# ─── Bottom methodology bar ───────────────────────────────────────────────
add_rounded_rect(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.45), CARD_BG, BORDER)
add_textbox(slide, Inches(1.0), Inches(6.85), Inches(11.3), Inches(0.35),
            "Methodology: Adapted Scrum — 2-week fixed-length sprints, backlog refinement, sprint review with supervisor, written retrospective per cycle",
            font_size=10, color=MUTED, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
            "1/1", font_size=10, color=MUTED, alignment=PP_ALIGN.RIGHT)

# ─── Save ──────────────────────────────────────────────────────────────────
output_path = "Sprint_Overview.pptx"
prs.save(output_path)
print(f"[OK] Saved: {output_path}")
